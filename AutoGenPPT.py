'''
from pptx import Presentation
import os
 
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[8])
placeholder = slide.placeholders[1]
picture = placeholder.insert_picture('GTA5_0.jpg')
prs.save("ESEMPIO.pptx")
os.startfile("ESEMPIO.pptx")
'''

##
##  Copy & Paste Tool for images to PowerPoint(.pptx)
##
import pptx
import pptx.util
import glob
import imageio
import os

OUTPUT_TAG = "BBB"

root_dir = "D:\BBB\Grade8Autumn\HighLevelCourses"
path = "D:\BBB\Grade8Autumn\HighLevelCourses\Module1_HowToLearnEnglish_Volcabulary"



bmp_file = path + "\*.bmp"
for root,dirs,files in os.walk(root_dir):
    for dir in dirs:
        # new
        prs = pptx.Presentation()

        # open
        # prs_exists = pptx.Presentation("some_presentation.pptx")

        # default slide width
        #prs.slide_width = 9144000
        # slide height @ 4:3
        #prs.slide_height = 6858000
        # slide height @ 16:9
        #prs.slide_height = 5143500

        # title slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        # blank slide
        #slide = prs.slides.add_slide(prs.slide_layouts[6])


        #pic_left  = int(prs.slide_width * 0.15)
        #pic_top   = int(prs.slide_height * 0.1)
        pic_left  = int(prs.slide_width * 0.0)
        pic_top   = int(prs.slide_height * 0.1)
        pic_width = int(prs.slide_width * 1.0)
        print(dir)
        bmp_file = root_dir + "\\" + dir + "\\*.bmp"
        #"\" + folders + "*.bmp"
        for g in glob.glob(bmp_file):
            pic_left  = int(prs.slide_width * 0.0)
            pic_top   = int(prs.slide_height * 0.1)
            pic_width = int(prs.slide_width * 1.0)
            print(g)
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            #tb = slide.shapes.add_textbox(0, 0, prs.slide_width, pic_top / 2)
            #p = tb.text_frame.add_paragraph()
            #p.text = g
            #p.font.size = pptx.util.Pt(14)

            img = imageio.imread(g)
            if(img.shape[0] < img.shape[1]):
                pic_height = int(pic_width * img.shape[0] / img.shape[1])
            else:
                pic_height = int(prs.slide_height * 0.8)
                pic_width = int(pic_height * img.shape[1] / img.shape[0])
            #pic   = slide.shapes.add_picture(g, pic_left, pic_top)
            pic   = slide.shapes.add_picture(g, pic_left, pic_top, pic_width, pic_height)

        OutFileName = OUTPUT_TAG + dir
        prs.save("%s.pptx" % OutFileName)